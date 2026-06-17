from pptx import Presentation
from pptx.util import Pt
import os


def generate_slides(script, topic, week):

    prs = Presentation()

    for scene in script["scenes"]:

        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Title
        slide.shapes.title.text = scene["slide_title"]

        # Body
        body = slide.placeholders[1].text_frame
        body.clear()

        for point in scene["bullet_points"]:
            p = body.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(22)

    output_folder = os.path.join(
        "courses",
        topic.replace(" ", "_"),
        f"Week_{week}",
        "slides"
    )

    os.makedirs(output_folder, exist_ok=True)

    filename = os.path.join(
        output_folder,
        "lesson.pptx"
    )

    prs.save(filename)

    return filename